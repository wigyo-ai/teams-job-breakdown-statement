from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

def rgb(r, g, b):
    return RGBColor(r, g, b)

# Colors from the image
C_NAVY       = rgb(30, 40, 60)      # Microsoft Teams / Azure Blob
C_PURPLE     = rgb(100, 80, 180)    # Webhook Service
C_TEAL       = rgb(30, 110, 110)    # Orchestrator
C_GREEN      = rgb(30, 120, 80)     # Document Generator
C_H2O_PURPLE = rgb(120, 80, 170)    # H2O box
C_ORANGE     = rgb(190, 90, 30)     # SharePoint box
C_BG         = rgb(245, 245, 248)   # Slide background
C_WHITE      = rgb(255, 255, 255)
C_ARROW_TEXT = rgb(60, 60, 80)
C_LIGHT_GRAY = rgb(220, 220, 230)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def fill_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

fill_slide_bg(slide, C_BG)

def add_rect(slide, x, y, w, h, fill_color, text_lines, font_sizes, bold_flags,
             text_color=C_WHITE, radius=None, port_text=None, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left  = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top   = Pt(4)
    tf.margin_bottom= Pt(4)

    for i, (line, fsize, bold) in enumerate(zip(text_lines, font_sizes, bold_flags)):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(fsize)
        run.font.bold = bold
        run.font.color.rgb = text_color
        run.font.name = "Calibri"

    # Port badge (top-right)
    if port_text:
        badge = slide.shapes.add_shape(
            1,
            Inches(x + w - 0.75), Inches(y + h/2 - 0.18),
            Inches(0.65), Inches(0.30)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(255, 255, 255)
        badge.fill.fore_color.theme_color  # noop
        badge.line.fill.background()
        btf = badge.text_frame
        btf.margin_left = Pt(2)
        btf.margin_right = Pt(2)
        btf.margin_top = Pt(1)
        p = btf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = port_text
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = fill_color
        run.font.name = "Calibri"

    return shape

def add_circle_badge(slide, x, y, r, fill_color, label, label_color=C_WHITE):
    d = r * 2
    shape = slide.shapes.add_shape(
        9,  # oval
        Inches(x - r), Inches(y - r), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = label_color
    run.font.name = "Calibri"

def add_text_box(slide, x, y, w, h, text, fsize=9, bold=False,
                 color=C_ARROW_TEXT, align=PP_ALIGN.CENTER, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fsize)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"

def add_arrow_label(slide, x, y, w, text):
    add_text_box(slide, x, y, w, 0.25,
                 f"\u21d5  {text}", fsize=8.5, color=rgb(80, 80, 110))

# ── Title ────────────────────────────────────────────────────────────────────
# Green circle icon
add_circle_badge(slide, 0.42, 0.32, 0.18, rgb(50, 180, 100), "")
add_text_box(slide, 0.60, 0.18, 4, 0.35, "Architecture Overview",
             fsize=18, bold=True, color=rgb(30, 30, 50), align=PP_ALIGN.LEFT)

# ── Layout constants ─────────────────────────────────────────────────────────
LX   = 0.25   # left column x
LW   = 8.55   # left column width
RX   = 9.05   # right column x
RW   = 4.0    # right column width

ROW_H = 0.62  # service box height
GAP   = 0.32  # gap between boxes (for arrow text)

Y0 = 0.65     # top of first box

# Row Y positions
y_teams  = Y0
y_wh     = y_teams + ROW_H + GAP
y_orch   = y_wh    + ROW_H + GAP
y_docgen = y_orch  + ROW_H + GAP
y_blob   = y_docgen + ROW_H + GAP

# ── Service boxes ─────────────────────────────────────────────────────────────

# 1. Microsoft Teams
add_rect(slide, LX, y_teams, LW, ROW_H, C_NAVY,
         ["Microsoft Teams (Azure Bot Service)",
          "External — user interface via Bot Framework REST API (HTTPS)"],
         [11, 8.5], [True, False])
# up-arrow icon
add_circle_badge(slide, LX + 0.28, y_teams + ROW_H/2, 0.17, rgb(60, 80, 130), "↑")

# Arrow 1
add_arrow_label(slide, LX + 0.3, y_teams + ROW_H + 0.04, LW - 0.5,
                "HTTPS POST /webhook/teams · RS256 JWT Bearer token")

# 2. Webhook Service
add_rect(slide, LX, y_wh, LW, ROW_H, C_PURPLE,
         ["Webhook Service — FastAPI",
          "Token validation · Activity parsing · Message normalization"],
         [11, 8.5], [True, False], port_text=":8000")
add_circle_badge(slide, LX + 0.28, y_wh + ROW_H/2, 0.17, rgb(80, 60, 160), "1",
                 label_color=C_WHITE)

# Arrow 2
add_arrow_label(slide, LX + 0.3, y_wh + ROW_H + 0.04, LW - 0.5,
                "HTTP POST /process · NormalisedMessage JSON (internal)")

# 3. Orchestrator Service
add_rect(slide, LX, y_orch, LW, ROW_H, C_TEAL,
         ["Orchestrator Service — FastAPI",
          "Session management · Phase 1 state machine · Phase 2 RAG hybrid"],
         [11, 8.5], [True, False], port_text=":8001")
add_circle_badge(slide, LX + 0.28, y_orch + ROW_H/2, 0.17, rgb(20, 90, 90), "2",
                 label_color=C_WHITE)

# Arrow 3
add_arrow_label(slide, LX + 0.3, y_orch + ROW_H + 0.04, LW - 0.5,
                "HTTP POST /generate · JBS JSON payload (internal)")

# 4. Document Generator Service
add_rect(slide, LX, y_docgen, LW, ROW_H, C_GREEN,
         ["Document Generator Service — FastAPI",
          "Template rendering · .docx assembly · Azure Blob upload · SAS URL"],
         [11, 8.5], [True, False], port_text=":8002")
add_circle_badge(slide, LX + 0.28, y_docgen + ROW_H/2, 0.17, rgb(20, 100, 60), "3",
                 label_color=C_WHITE)

# Arrow 4
add_arrow_label(slide, LX + 0.3, y_docgen + ROW_H + 0.04, LW - 0.5,
                "Azure Blob Storage SDK · SAS token (15-min read-only)")

# 5. Azure Blob Storage
add_rect(slide, LX, y_blob, LW, ROW_H, C_NAVY,
         ["Azure Blob Storage",
          "Container: jbs-documents · Prefix: jbs-documents/"],
         [11, 8.5], [True, False])
add_circle_badge(slide, LX + 0.28, y_blob + ROW_H/2, 0.17, rgb(40, 60, 110), "↓")

# ── Right column boxes ────────────────────────────────────────────────────────
rh1 = 1.65   # H2O box height
rh2 = 1.55   # SharePoint box height
ry1 = y_wh   # align with Webhook row

add_rect(slide, RX, ry1, RW, rh1, C_H2O_PURPLE,
         ["H2O Enterprise h2oGPTe", "RAG Engine",
          "5 site-category collections", "Vector DB: hnswlib",
          "Auth: API Key (Key Vault)"],
         [12, 9, 9, 9, 9],
         [True, False, False, False, False])

# "called by Orchestrator" label
add_text_box(slide, RX, ry1 + rh1 + 0.02, RW, 0.22,
             "called by Orchestrator", fsize=8, italic=True,
             color=rgb(120, 120, 140), align=PP_ALIGN.CENTER)

ry2 = ry1 + rh1 + 0.30
add_rect(slide, RX, ry2, RW, rh2, C_ORANGE,
         ["SharePoint Online", "Knowledge Base source",
          "5 document libraries", "Ingested via Graph API",
          "Daily sync: 02:00 UTC"],
         [12, 9, 9, 9, 9],
         [True, False, False, False, False])

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "/Users/lmccoy/job-breakdown-statement02/architecture_overview.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
